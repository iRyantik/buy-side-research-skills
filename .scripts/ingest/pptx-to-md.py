#!/usr/bin/env python3
"""pptx-to-md.py — Convert PowerPoint deck to markdown.

Usage:
    python pptx-to-md.py <input.pptx> [--output <path>]

One ## Slide N section per slide, with speaker notes where available.
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path


def convert(input_path: Path, output_path: Path) -> dict:
    from pptx import Presentation

    deck = Presentation(str(input_path))
    parts = [f"# {input_path.name}", ""]
    slide_count = 0

    for idx, slide in enumerate(deck.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

        if not slide_texts:
            continue

        slide_count += 1
        parts.append(f"## Slide {idx}")
        parts.append("")
        for text in slide_texts:
            parts.append(text)
            parts.append("")

        # Speaker notes
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append("### Speaker Notes")
                parts.append("")
                parts.append(notes)
                parts.append("")
        except Exception:
            pass

    output_path.parent.mkdir(parents=True, exist_ok=True)
    text = "\n".join(parts)
    output_path.write_text(text, encoding="utf-8")

    return {
        "status": "converted",
        "engine": "python-pptx",
        "slides_total": len(deck.slides),
        "slides_with_text": slide_count,
        "chars": len(text),
    }


def main():
    parser = argparse.ArgumentParser(description="Convert PowerPoint to markdown")
    parser.add_argument("input", help="PPTX file path")
    parser.add_argument("--output", "-o", help="Output markdown path")
    args = parser.parse_args()

    src = Path(args.input)
    if not src.exists():
        print(f"ERROR: file not found: {src}", file=sys.stderr)
        sys.exit(1)

    dst = Path(args.output) if args.output else src.with_suffix(".md")
    result = convert(src, dst)
    print(f"OK {result['slides_with_text']}/{result['slides_total']} slides, {result['chars']:,} chars → {dst}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
