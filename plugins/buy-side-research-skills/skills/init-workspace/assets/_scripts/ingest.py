#!/usr/bin/env python3
"""Convert raw research files into source-tracked markdown cache files."""

from __future__ import annotations

import argparse
import csv
import datetime as dt
from dataclasses import dataclass
import hashlib
import html
import importlib.util
import json
import os
from pathlib import Path
import re
import shutil
import sys
from typing import Any

from PIL import Image


SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".xlsx",
    ".xlsm",
    ".xls",
    ".docx",
    ".pptx",
    ".pdf",
}

SEC_FORM_PATTERN = re.compile(r"\b(10[-_ ]?K|10[-_ ]?Q|8[-_ ]?K|20[-_ ]?F|40[-_ ]?F)\b", re.IGNORECASE)
SEC_TEXT_PATTERNS = (
    "UNITED STATES SECURITIES AND EXCHANGE COMMISSION",
    "<SEC-DOCUMENT",
    "EDGAR",
    "FORM 10-K",
    "FORM 10-Q",
    "FORM 8-K",
    "FORM 20-F",
)
NUMBER_PATTERN = re.compile(r"(?<![A-Za-z])[-+]?\$?\d[\d,]*(?:\.\d+)?%?")


class IngestError(RuntimeError):
    pass


@dataclass
class DocumentProfile:
    extension: str
    document_type: str
    page_count: int | None = None
    table_count: int | None = None
    text_chars: int | None = None
    ocr_required: bool = False
    sec_filing: bool = False


@dataclass
class FigureEntry:
    """Metadata for one extracted figure/chart."""
    figure_id: str
    page: int | None = None
    bbox: tuple[float, float, float, float] | None = None
    caption: str | None = None
    classification: str | None = None
    classification_confidence: float | None = None
    image_rel_path: str | None = None
    width_px: int | None = None
    height_px: int | None = None
    description: str | None = None


@dataclass
class FigureManifest:
    """Manifest of all figures extracted from one source document."""
    source_stem: str
    figures: list[FigureEntry]
    images_dir: str


@dataclass
class ConversionResult:
    markdown: str
    converter: str
    precision: str
    precision_level: str
    document_type: str
    route: str
    page_count: int | None = None
    table_count: int | None = None
    ocr_required: bool = False
    figure_manifest: FigureManifest | None = None
    dependency_status: dict[str, Any] | None = None


def utc_now() -> str:
    return dt.datetime.now(dt.timezone.utc).replace(microsecond=0).isoformat()


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def read_text_lossy(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def module_available(module_name: str) -> bool:
    return importlib.util.find_spec(module_name) is not None


def dependency_matrix() -> dict[str, Any]:
    packages = {
        "docling": {
            "available": module_available("docling"),
            "module": "docling",
            "install_hint": "pip install docling",
        },
        "edgartools": {
            "available": module_available("edgar"),
            "module": "edgar",
            "install_hint": "pip install edgartools",
        },
        "pymupdf4llm": {
            "available": module_available("pymupdf4llm"),
            "module": "pymupdf4llm",
            "install_hint": "pip install pymupdf4llm",
        },
        "openpyxl": {
            "available": module_available("openpyxl"),
            "module": "openpyxl",
            "install_hint": "pip install openpyxl",
        },
        "python-pptx": {
            "available": module_available("pptx"),
            "module": "pptx",
            "install_hint": "pip install python-pptx",
        },
        "python-docx": {
            "available": module_available("docx"),
            "module": "docx",
            "install_hint": "pip install python-docx",
        },
        "pdfplumber": {
            "available": module_available("pdfplumber"),
            "module": "pdfplumber",
            "install_hint": "pip install pdfplumber",
        },
        "pypdf": {
            "available": module_available("pypdf"),
            "module": "pypdf",
            "install_hint": "pip install pypdf",
        },
        "Pillow": {
            "available": module_available("PIL"),
            "module": "PIL",
            "install_hint": "pip install Pillow",
        },
    }
    return {
        "python": {
            "version": sys.version.split()[0],
            "executable": sys.executable,
        },
        "packages": packages,
        "binaries": {},
        "env": {
            "EDGAR_IDENTITY": {
                "configured": bool(os.getenv("EDGAR_IDENTITY")),
                "value": os.getenv("EDGAR_IDENTITY"),
            }
        },
    }


def dependency_snapshot(keys: list[str]) -> dict[str, Any]:
    matrix = dependency_matrix()
    packages = matrix["packages"]
    return {
        "packages": {key: packages.get(key) for key in keys if key in packages},
        "binaries": matrix["binaries"],
        "env": matrix["env"],
    }


def require_package(package_key: str, module_name: str | None = None) -> None:
    module_name = module_name or package_key
    if not module_available(module_name):
        matrix = dependency_matrix()
        hint = matrix["packages"].get(package_key, {}).get("install_hint", f"pip install {package_key}")
        raise IngestError(f"Missing optional dependency: {package_key}. Run bootstrap-ingest-deps.ps1 or `{hint}`.")


def inspect_pdf(path: Path) -> tuple[int | None, int | None, int | None]:
    page_count: int | None = None
    text_chars: int | None = None
    table_count: int | None = None

    if module_available("pypdf"):
        try:
            from pypdf import PdfReader  # type: ignore

            reader = PdfReader(str(path))
            page_count = len(reader.pages)
            texts = []
            for page in reader.pages[:20]:
                texts.append(page.extract_text() or "")
            text_chars = sum(len(text.strip()) for text in texts)
        except Exception:
            pass

    if module_available("pdfplumber"):
        try:
            import pdfplumber  # type: ignore

            count = 0
            with pdfplumber.open(str(path)) as pdf:
                if page_count is None:
                    page_count = len(pdf.pages)
                for page in pdf.pages:
                    count += len(page.extract_tables() or [])
            table_count = count
        except Exception:
            pass

    return page_count, text_chars, table_count


def extract_pdf_text_probe(path: Path, max_chars: int = 8000) -> str:
    if not module_available("pypdf"):
        return ""
    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(str(path))
        chunks = []
        for page in reader.pages[:5]:
            chunks.append(page.extract_text() or "")
            if sum(len(chunk) for chunk in chunks) >= max_chars:
                break
        return "\n".join(chunks)[:max_chars]
    except Exception:
        return ""


def detect_format(path: Path) -> DocumentProfile:
    extension = path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise IngestError(f"Unsupported file extension: {extension}")

    if extension == ".pdf":
        page_count, text_chars, table_count = inspect_pdf(path)
        probe_text = extract_pdf_text_probe(path)
        sec_filing = bool(SEC_FORM_PATTERN.search(path.name)) or any(pattern in probe_text.upper() for pattern in SEC_TEXT_PATTERNS)
        ocr_required = bool(page_count and (text_chars is None or text_chars < 20))
        return DocumentProfile(
            extension=extension,
            document_type="sec-filing-pdf" if sec_filing else "pdf",
            page_count=page_count,
            table_count=table_count,
            text_chars=text_chars,
            ocr_required=ocr_required,
            sec_filing=sec_filing,
        )

    type_map = {
        ".txt": "text",
        ".md": "markdown",
        ".markdown": "markdown",
        ".csv": "csv",
        ".xlsx": "excel-workbook",
        ".xlsm": "excel-workbook",
        ".xls": "legacy-excel-workbook",
        ".docx": "word-document",
        ".pptx": "powerpoint-deck",
    }
    return DocumentProfile(extension=extension, document_type=type_map[extension])


def markdown_header(source: Path, result: ConversionResult) -> str:
    modified = dt.datetime.fromtimestamp(source.stat().st_mtime, dt.timezone.utc)
    metadata = {
        "source_path": str(source),
        "source_sha256": sha256_file(source),
        "source_modified_utc": modified.replace(microsecond=0).isoformat(),
        "converter": result.converter,
        "converted_at_utc": utc_now(),
        "precision": result.precision,
        "precision_level": result.precision_level,
        "document_type": result.document_type,
        "route": result.route,
        "page_count": result.page_count if result.page_count is not None else "unknown",
        "table_count": result.table_count if result.table_count is not None else "unknown",
        "ocr_required": str(result.ocr_required).lower(),
        "figures_extracted": str(len(result.figure_manifest.figures)) if result.figure_manifest else "0",
        "figures_described": str(sum(1 for f in result.figure_manifest.figures if f.description)) if result.figure_manifest else "0",
    }
    body = "\n".join(f"{key}: {value}" for key, value in metadata.items())
    return f"<!--\n{body}\n-->\n\n"


def render_figure_block(fig: FigureEntry, source_stem: str) -> str:
    """Generate a structured markdown figure block."""
    parts = [
        f"<!-- FIGURE {fig.figure_id} | page: {fig.page or '?'} | classification: {fig.classification or 'unknown'} | confidence: {fig.classification_confidence or '0'} -->",
        f"![Figure: {fig.caption or 'Uncaptioned figure'}]({fig.image_rel_path})",
        f"*Page {fig.page or '?'} -- {fig.classification or 'unknown'} ({fig.classification_confidence or 0})*",
    ]
    if fig.description:
        parts.append(f"> **Figure Description:** {fig.description}")
    else:
        parts.append("> [Figure description pending]")
    parts.append("<!-- /FIGURE -->")
    return "\n\n" + "\n".join(parts) + "\n"


def convert_markdown(path: Path, profile: DocumentProfile) -> ConversionResult:
    return ConversionResult(
        markdown=read_text_lossy(path),
        converter="native-markdown",
        precision="native markdown; verify quoted facts against original source",
        precision_level="native",
        document_type=profile.document_type,
        route="native-markdown",
        dependency_status=dependency_snapshot([]),
    )


def convert_text(path: Path, profile: DocumentProfile) -> ConversionResult:
    text = read_text_lossy(path)
    return ConversionResult(
        markdown=f"# {path.name}\n\n```text\n{text}\n```\n",
        converter="native-text",
        precision="plain text extraction",
        precision_level="native",
        document_type=profile.document_type,
        route="native-text",
        dependency_status=dependency_snapshot([]),
    )


def convert_csv(path: Path, profile: DocumentProfile) -> ConversionResult:
    text = read_text_lossy(path)
    rows = list(csv.reader(text.splitlines()))
    preview_rows = rows[:200]
    output = [f"# {path.name}", "", "## CSV Preview", ""]
    if preview_rows:
        width = max(len(row) for row in preview_rows)
        normalized = [row + [""] * (width - len(row)) for row in preview_rows]
        output.append("| " + " | ".join(html.escape(cell) for cell in normalized[0]) + " |")
        output.append("| " + " | ".join("---" for _ in normalized[0]) + " |")
        for row in normalized[1:]:
            output.append("| " + " | ".join(html.escape(cell) for cell in row) + " |")
    else:
        output.append("[empty csv]")
    output.extend(["", "## Raw CSV", "", "```csv", text, "```", ""])
    return ConversionResult(
        markdown="\n".join(output),
        converter="native-csv",
        precision="CSV parsed with Python csv; inspect original for dialect/encoding issues",
        precision_level="native",
        document_type=profile.document_type,
        route="native-csv",
        dependency_status=dependency_snapshot([]),
    )


def convert_xlsx(path: Path, profile: DocumentProfile) -> ConversionResult:
    require_package("openpyxl", "openpyxl")
    scripts_dir = Path(__file__).resolve().parent
    if str(scripts_dir) not in sys.path:
        sys.path.insert(0, str(scripts_dir))
    from ingest_xlsx import workbook_to_markdown

    return ConversionResult(
        markdown=workbook_to_markdown(path),
        converter="openpyxl",
        precision="workbook structure extraction; verify formulas and key numbers in Excel",
        precision_level="structured",
        document_type=profile.document_type,
        route="excel-openpyxl-dual-load",
        dependency_status=dependency_snapshot(["openpyxl"]),
    )


def convert_docling(
    path: Path,
    profile: DocumentProfile,
    route: str,
    cache_parent: Path | None = None,
    extract_figures: bool = True,
) -> ConversionResult:
    require_package("docling", "docling")
    from docling.document_converter import DocumentConverter  # type: ignore

    converter = DocumentConverter()
    try:
        pdf_opt = converter.format_to_options.get("pdf")
        if pdf_opt and hasattr(pdf_opt, "pipeline_options"):
            pdf_opt.pipeline_options.generate_picture_images = True
            pdf_opt.pipeline_options.images_scale = 1.5
    except Exception:
        pass

    result = converter.convert(str(path))
    markdown = result.document.export_to_markdown()

    figure_manifest = None
    if extract_figures and cache_parent is not None:
        try:
            source_stem = re.sub(r'[^\w\-]', '_', path.stem)[:80]
            figure_manifest = extract_docling_figures(result, source_stem, path, cache_parent)
            if figure_manifest and figure_manifest.figures:
                # Append figure blocks to markdown
                markdown += "\n\n## Extracted Figures\n\n"
                for fig in figure_manifest.figures:
                    markdown += render_figure_block(fig, source_stem)
        except Exception:
            pass

    return ConversionResult(
        markdown=markdown,
        converter="docling",
        precision="Docling markdown conversion; verify financial tables and layout-sensitive content against original",
        precision_level="structured",
        document_type=profile.document_type,
        route=route,
        page_count=profile.page_count,
        table_count=profile.table_count,
        ocr_required=profile.ocr_required,
        figure_manifest=figure_manifest,
        dependency_status=dependency_snapshot(["docling", "pypdf", "pdfplumber", "edgartools"]),
    )


def convert_pymupdf4llm(
    path: Path,
    profile: DocumentProfile,
    route: str,
    cache_parent: Path | None = None,
    extract_figures: bool = True,
) -> ConversionResult:
    require_package("pymupdf4llm", "pymupdf4llm")
    import pymupdf4llm  # type: ignore

    markdown = pymupdf4llm.to_markdown(str(path))

    figure_manifest = None
    if extract_figures and cache_parent is not None:
        try:
            source_stem = re.sub(r'[^\w\-]', '_', path.stem)[:80]
            figure_manifest = extract_pymupdf_figures(path, source_stem, cache_parent)
            if figure_manifest and figure_manifest.figures:
                markdown += "\n\n## Extracted Figures\n\n"
                for fig in figure_manifest.figures:
                    markdown += render_figure_block(fig, source_stem)
        except Exception:
            pass

    return ConversionResult(
        markdown=markdown,
        converter="pymupdf4llm",
        precision="PyMuPDF4LLM text extraction; fast CPU path suitable for text-heavy documents, verify tables against original",
        precision_level="structured",
        document_type=profile.document_type,
        route=route,
        page_count=profile.page_count,
        table_count=profile.table_count,
        ocr_required=profile.ocr_required,
        figure_manifest=figure_manifest,
        dependency_status=dependency_snapshot(["pymupdf4llm", "pypdf", "pdfplumber"]),
    )


def convert_docx_fallback(path: Path, profile: DocumentProfile) -> ConversionResult:
    require_package("python-docx", "docx")
    import docx  # type: ignore

    document = docx.Document(str(path))
    parts = [f"# {path.name}", ""]
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
            parts.append("")
    for idx, table in enumerate(document.tables, start=1):
        parts.append(f"## Table {idx}")
        for row in table.rows:
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            parts.append("| " + " | ".join(html.escape(cell) for cell in cells) + " |")
        parts.append("")
    return ConversionResult(
        markdown="\n".join(parts),
        converter="python-docx",
        precision="DOCX fallback text/table extraction; verify layout-sensitive content in original",
        precision_level="degraded",
        document_type=profile.document_type,
        route="docx-python-docx-fallback",
        table_count=len(document.tables),
        dependency_status=dependency_snapshot(["python-docx"]),
    )


def convert_pptx_fallback(
    path: Path,
    profile: DocumentProfile,
    cache_parent: Path | None = None,
    extract_figures: bool = True,
) -> ConversionResult:
    require_package("python-pptx", "pptx")
    from pptx import Presentation  # type: ignore

    deck = Presentation(str(path))
    parts = [f"# {path.name}", ""]
    for idx, slide in enumerate(deck.slides, start=1):
        parts.append(f"## Slide {idx}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    parts.append(text)
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.extend(["", "### Speaker Notes", notes])
        except Exception:
            pass
        parts.append("")

    figure_manifest = None
    if extract_figures and cache_parent is not None:
        try:
            source_stem = re.sub(r'[^\w\-]', '_', path.stem)[:80]
            figure_manifest = extract_pptx_figures(path, source_stem, cache_parent)
            if figure_manifest and figure_manifest.figures:
                parts.append("\n\n## Extracted Figures\n\n")
                for fig in figure_manifest.figures:
                    parts.append(render_figure_block(fig, source_stem))
        except Exception:
            pass

    return ConversionResult(
        markdown="\n".join(parts),
        converter="python-pptx",
        precision="PPTX fallback text/notes extraction; verify charts/images in original deck",
        precision_level="degraded",
        document_type=profile.document_type,
        route="pptx-python-pptx-fallback",
        page_count=len(deck.slides),
        figure_manifest=figure_manifest,
        dependency_status=dependency_snapshot(["python-pptx"]),
    )


def convert_pdf_pypdf_fallback(path: Path, profile: DocumentProfile, route: str) -> ConversionResult:
    require_package("pypdf", "pypdf")
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(str(path))
    parts = [f"# {path.name}", ""]
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        parts.extend([f"## Page {idx}", "", text.strip() or "[no extractable text]", ""])
    return ConversionResult(
        markdown="\n".join(parts),
        converter="pypdf",
        precision="PDF text-layer fallback only; tables, scans, and layout require manual verification",
        precision_level="degraded",
        document_type=profile.document_type,
        route=route,
        page_count=profile.page_count or len(reader.pages),
        table_count=profile.table_count,
        ocr_required=profile.ocr_required,
        dependency_status=dependency_snapshot(["pypdf", "pdfplumber"]),
    )


def ensure_sec_route_ready(profile: DocumentProfile) -> None:
    if not profile.sec_filing:
        return
    require_package("edgartools", "edgar")
    if not os.getenv("EDGAR_IDENTITY"):
        raise IngestError(
            "SEC filing detected but EDGAR_IDENTITY is not configured. Run bootstrap-ingest-deps.ps1 -EdgarIdentity \"Name email@domain.com\"."
        )


def convert_pdf(
    path: Path,
    profile: DocumentProfile,
    cache_parent: Path | None = None,
    extract_figures: bool = True,
) -> ConversionResult:
    if profile.ocr_required:
        if module_available("docling"):
            result = convert_docling(path, profile, route="pdf-docling-scanned", cache_parent=cache_parent, extract_figures=extract_figures)
            result.precision = "Scanned PDF via Docling; text may contain OCR errors, verify key numbers and tables against original scan. For critical scanned documents, prefer Claude Vision review."
            return result
        if module_available("pymupdf4llm"):
            result = convert_pymupdf4llm(path, profile, route="pdf-pymupdf4llm-scanned-fallback", cache_parent=cache_parent, extract_figures=extract_figures)
            result.precision = "Scanned PDF via PyMuPDF4LLM fallback; text may contain OCR errors, verify key numbers against original. For critical scanned documents, prefer Claude Vision review."
            return result
        raise IngestError(
            "Scanned PDF detected but neither docling nor pymupdf4llm is available. "
            "Run bootstrap-ingest-deps.ps1; no cache was written."
        )

    if profile.sec_filing:
        ensure_sec_route_ready(profile)
        try:
            result = convert_docling(path, profile, route="sec-filing-edgartools-xbrl-docling-narrative", cache_parent=cache_parent, extract_figures=extract_figures)
            note = (
                "\n\n## SEC XBRL Route Note\n\n"
                "EdgarTools dependency and EDGAR_IDENTITY are available. This cache converts the local filing narrative; "
                "financial statement numbers should still be reconciled to SEC/XBRL source before use.\n"
            )
            result.markdown = result.markdown.rstrip() + note
            return result
        except IngestError:
            raise
        except Exception as exc:
            raise IngestError(f"SEC filing route failed before cache write: {exc}") from exc

    table_heavy = bool(profile.table_count and profile.table_count >= 15)

    if table_heavy and module_available("docling"):
        try:
            return convert_docling(path, profile, route="pdf-docling-table-heavy", cache_parent=cache_parent, extract_figures=extract_figures)
        except Exception as exc:
            if module_available("pymupdf4llm"):
                fallback = convert_pymupdf4llm(path, profile, route="pdf-pymupdf4llm-fallback-after-docling-error", cache_parent=cache_parent, extract_figures=extract_figures)
                fallback.precision = f"{fallback.precision}; Docling failed: {exc}"
                return fallback
            raise IngestError(f"Docling PDF conversion failed for table-heavy document: {exc}") from exc

    if module_available("pymupdf4llm"):
        try:
            return convert_pymupdf4llm(path, profile, route="pdf-pymupdf4llm-text", cache_parent=cache_parent, extract_figures=extract_figures)
        except Exception as exc:
            if module_available("docling"):
                fallback = convert_docling(path, profile, route="pdf-docling-fallback-after-pymupdf4llm-error", cache_parent=cache_parent, extract_figures=extract_figures)
                fallback.precision = f"{fallback.precision}; PyMuPDF4LLM failed: {exc}"
                return fallback
            raise IngestError(f"PyMuPDF4LLM conversion failed: {exc}") from exc

    if module_available("docling"):
        return convert_docling(path, profile, route="pdf-docling", cache_parent=cache_parent, extract_figures=extract_figures)

    if profile.text_chars and profile.text_chars >= 20 and module_available("pypdf"):
        return convert_pdf_pypdf_fallback(path, profile, route="pdf-pypdf-fallback-no-docling-no-pymupdf4llm")

    raise IngestError("PDF conversion requires Docling or PyMuPDF4LLM. Run bootstrap-ingest-deps.ps1; no cache was written.")


def extract_docling_figures(
    docling_result,
    source_stem: str,
    source_path: Path,
    cache_parent: Path,
) -> FigureManifest:
    """Extract embedded picture items from a Docling conversion result."""
    images_dir = f"_images/{source_stem}/"
    abs_images_dir = cache_parent / "_images" / source_stem
    abs_images_dir.mkdir(parents=True, exist_ok=True)

    figures: list[FigureEntry] = []
    content_hashes: set[str] = set()
    fig_counter = 0

    try:
        from docling.datamodel.base_models import Size
        from docling.datamodel.pipeline_options import PdfPipelineOptions
    except ImportError:
        pass

    for item, _depth in docling_result.document.iterate_items():
        if type(item).__name__ != 'PictureItem':
            continue
        page = item.prov[0].page_no if item.prov else None
        bbox = None
        try:
            if item.prov and item.prov[0].bbox:
                b = item.prov[0].bbox
                bbox = (b.l, b.t, b.r, b.b)
        except Exception:
            pass

        # Try to get image from the picture item
        img = None
        if hasattr(item, "image") and item.image is not None:
            try:
                img = item.image
            except Exception:
                pass

        if img is None:
            continue

        # Convert to PIL Image
        try:
            import io
            import hashlib
            from docling_core.types.doc.document import ImageRef

            if isinstance(img, ImageRef):
                pil_img = img.pil_image
                if pil_img is None:
                    continue
            elif isinstance(img, (bytes, bytearray)):
                pil_img = Image.open(io.BytesIO(img))
            elif hasattr(img, 'pil_image'):
                pil_img = img.pil_image
                if pil_img is None:
                    continue
            else:
                pil_img = img

            if pil_img is None:
                continue
        except Exception:
            continue

        w, h = pil_img.size
        # Filter out tiny images (icons, bullets, decorative elements)
        if max(w, h) < 200:
            continue

        # Deduplicate by image content hash
        import io as _io
        buf = _io.BytesIO()
        pil_img.save(buf, format="PNG")
        content_hash = hashlib.sha256(buf.getvalue()).hexdigest()
        if content_hash in content_hashes:
            continue
        content_hashes.add(content_hash)

        fig_counter += 1
        figure_id = f"fig-{fig_counter:03d}"
        safe_stem = re.sub(r'[^\w\-]', '_', source_stem)[:80]
        filename = f"{figure_id}-p{page or 0}.png"
        filepath = abs_images_dir / filename
        # Resize to max 2048px on longest side
        if max(w, h) > 2048:
            ratio = 2048 / max(w, h)
            pil_img = pil_img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)
        pil_img.save(filepath, "PNG", optimize=True)
        rel_path = f"{images_dir}{filename}"

        # Try to get caption
        caption = None
        try:
            if hasattr(item, "captions") and item.captions:
                caption = " ".join(c.text for c in item.captions if hasattr(c, "text"))
        except Exception:
            pass

        # Try to get classification
        classification = None
        confidence = None
        try:
            if hasattr(item, "annotations") and item.annotations:
                for ann in item.annotations:
                    cls_name = getattr(ann, "name", None) or getattr(ann, "label", None)
                    if cls_name:
                        classification = str(cls_name)
                        confidence = getattr(ann, "confidence", None)
                        break
        except Exception:
            pass

        figures.append(FigureEntry(
            figure_id=figure_id,
            page=page,
            bbox=bbox,
            caption=caption,
            classification=classification,
            classification_confidence=confidence,
            image_rel_path=rel_path,
            width_px=pil_img.width,
            height_px=pil_img.height,
            description=None,
        ))

    return FigureManifest(source_stem=source_stem, figures=figures, images_dir=images_dir)


def extract_pymupdf_figures(
    source_path: Path,
    source_stem: str,
    cache_parent: Path,
) -> FigureManifest:
    """Extract images from PDF using PyMuPDF (fitz) for pymupdf4llm route."""
    images_dir = f"_images/{source_stem}/"
    abs_images_dir = cache_parent / "_images" / source_stem
    abs_images_dir.mkdir(parents=True, exist_ok=True)

    figures: list[FigureEntry] = []
    content_hashes: set[str] = set()
    fig_counter = 0

    try:
        import fitz
    except ImportError:
        return FigureManifest(source_stem=source_stem, figures=[], images_dir=images_dir)

    try:
        doc = fitz.open(str(source_path))
    except Exception:
        return FigureManifest(source_stem=source_stem, figures=[], images_dir=images_dir)

    import hashlib

    for page_num in range(len(doc)):
        page = doc[page_num]
        image_list = page.get_images(full=True)
        for img_info in image_list:
            xref = img_info[0]
            try:
                base = doc.extract_image(xref)
            except Exception:
                continue
            img_bytes = base.get("image")
            if not img_bytes:
                continue
            try:
                import io
                pil_img = Image.open(io.BytesIO(img_bytes))
            except Exception:
                continue
            w, h = pil_img.size
            # Stricter filter: fitz catches all layout images including header/footer
            if max(w, h) < 400:
                continue

            # Deduplicate by content hash
            content_hash = hashlib.sha256(img_bytes).hexdigest()
            if content_hash in content_hashes:
                continue
            content_hashes.add(content_hash)

            fig_counter += 1
            figure_id = f"fig-{fig_counter:03d}"
            safe_stem = re.sub(r'[^\w\-]', '_', source_stem)[:80]
            filename = f"{figure_id}-p{page_num + 1}.png"
            filepath = abs_images_dir / filename
            if max(w, h) > 2048:
                ratio = 2048 / max(w, h)
                pil_img = pil_img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)
            pil_img.save(filepath, "PNG", optimize=True)
            rel_path = f"{images_dir}{filename}"

            figures.append(FigureEntry(
                figure_id=figure_id,
                page=page_num + 1,
                image_rel_path=rel_path,
                width_px=pil_img.width,
                height_px=pil_img.height,
                description=None,
            ))

    doc.close()
    return FigureManifest(source_stem=source_stem, figures=figures, images_dir=images_dir)


def extract_pptx_figures(
    source_path: Path,
    source_stem: str,
    cache_parent: Path,
) -> FigureManifest:
    """Extract embedded images from PPTX slides."""
    images_dir = f"_images/{source_stem}/"
    abs_images_dir = cache_parent / "_images" / source_stem
    abs_images_dir.mkdir(parents=True, exist_ok=True)

    figures: list[FigureEntry] = []
    fig_counter = 0

    try:
        from pptx import Presentation
    except ImportError:
        return FigureManifest(source_stem=source_stem, figures=[], images_dir=images_dir)

    try:
        deck = Presentation(str(source_path))
    except Exception:
        return FigureManifest(source_stem=source_stem, figures=[], images_dir=images_dir)

    for slide_idx, slide in enumerate(deck.slides, start=1):
        for shape in slide.shapes:
            img_blob = None
            try:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    img_blob = shape.image.blob
                elif hasattr(shape, "image"):
                    img_blob = shape.image.blob
            except Exception:
                continue
            if not img_blob:
                continue
            try:
                import io
                pil_img = Image.open(io.BytesIO(img_blob))
            except Exception:
                continue
            w, h = pil_img.size
            if max(w, h) < 200:
                continue

            fig_counter += 1
            figure_id = f"fig-{fig_counter:03d}"
            safe_stem = re.sub(r'[^\w\-]', '_', source_stem)[:80]
            filename = f"{figure_id}-p{slide_idx}.png"
            filepath = abs_images_dir / filename
            if max(w, h) > 2048:
                ratio = 2048 / max(w, h)
                pil_img = pil_img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)
            pil_img.save(filepath, "PNG", optimize=True)
            rel_path = f"{images_dir}{filename}"

            figures.append(FigureEntry(
                figure_id=figure_id,
                page=slide_idx,
                image_rel_path=rel_path,
                width_px=pil_img.width,
                height_px=pil_img.height,
                description=None,
            ))

    return FigureManifest(source_stem=source_stem, figures=figures, images_dir=images_dir)


def route_converter(
    path: Path,
    profile: DocumentProfile,
    cache_parent: Path | None = None,
    extract_figures: bool = True,
) -> ConversionResult:
    if profile.extension in {".md", ".markdown"}:
        return convert_markdown(path, profile)
    if profile.extension == ".txt":
        return convert_text(path, profile)
    if profile.extension == ".csv":
        return convert_csv(path, profile)
    if profile.extension in {".xlsx", ".xlsm"}:
        return convert_xlsx(path, profile)
    if profile.extension == ".xls":
        try:
            return convert_xlsx(path, profile)
        except Exception:
            raise IngestError(
                "Legacy .xls format is not supported. Convert the file to .xlsx in Excel and re-ingest."
            )
    if profile.extension == ".docx":
        if module_available("docling"):
            return convert_docling(path, profile, route="docx-docling", cache_parent=cache_parent, extract_figures=extract_figures)
        return convert_docx_fallback(path, profile)
    if profile.extension == ".pptx":
        if module_available("docling"):
            try:
                docling_result = convert_docling(path, profile, route="pptx-docling", cache_parent=cache_parent, extract_figures=extract_figures)
                if module_available("pptx"):
                    notes_result = convert_pptx_fallback(path, profile, cache_parent=cache_parent, extract_figures=extract_figures)
                    docling_result.markdown = (
                        docling_result.markdown.rstrip()
                        + "\n\n## python-pptx Speaker Notes Fallback\n\n"
                        + notes_result.markdown
                    )
                return docling_result
            except Exception:
                return convert_pptx_fallback(path, profile, cache_parent=cache_parent, extract_figures=extract_figures)
        return convert_pptx_fallback(path, profile, cache_parent=cache_parent, extract_figures=extract_figures)
    if profile.extension == ".pdf":
        return convert_pdf(path, profile, cache_parent=cache_parent, extract_figures=extract_figures)
    raise IngestError(f"Unsupported file extension: {profile.extension}")


RAW_CATEGORIES = ("filings", "transcripts", "sellside", "industry", "irdecks", "datasets")
TOPIC_NAMESPACES = ("company", "industry", "theme", "pair")
TOPIC_OPERATION_DIRS = ("_raw", "_inbox", "_cache")

CATEGORY_FILENAME_HINTS = {
    "filings": (r"\b(10[-_]?K|10[-_]?Q|8[-_]?K|20[-_]?F|40[-_]?F|annual.report|proxy|prospectus)\b",),
    "transcripts": (r"\b(transcript|earnings.call|conference.call|investor.call)\b",),
    "sellside": (r"\b(initiat|rating|target.price|research.report|upgrade|downgrade)\b",),
    "irdecks": (r"\b(deck|presentation|investor.day|ir.deck|roadshow)\b",),
    "industry": (r"\b(industry|market.report|outlook|forecast|white.paper)\b",),
    "datasets": (),
}


def infer_category(path: Path, profile) -> str:
    fname = path.name.lower()
    for cat, patterns in CATEGORY_FILENAME_HINTS.items():
        for pat in patterns:
            if re.search(pat, fname):
                return cat
    ext = path.suffix.lower()
    if ext in (".xlsx", ".xlsm", ".xls", ".csv"):
        return "datasets"
    if getattr(profile, "sec_filing", False):
        return "filings"
    return "unclassified"


def discover_workspace(source: Path) -> Path:
    candidates = [source if source.is_dir() else source.parent, Path.cwd()]
    for candidate in candidates:
        for parent in [candidate, *candidate.parents]:
            if (parent / "industry").is_dir():
                return parent
    raise IngestError("Could not discover research workspace. Pass --workspace or run init first.")


def normalize_topic_arg(topic: str | None) -> str | None:
    if not topic:
        return None
    normalized = topic.replace("\\", "/").strip().strip("/")
    if normalized.startswith("topics/") or normalized.startswith("industry/"):
        # 'industry/pcb-equipment/companies/mycronic' -> keep as-is (it's already normalized)
        # 'topics/industry/pcb-equipment' -> strip 'topics/' prefix
        if normalized.startswith("topics/"):
            normalized = normalized[len("topics/") :]
    elif "/" not in normalized:
        # Bare slug like "optical-module-equipment" → assume industry/ structure
        normalized = f"industry/{normalized}"
    return normalized or None


def resolve_topic(source: Path, workspace: Path, explicit_topic: str | None) -> str:
    normalized_topic = normalize_topic_arg(explicit_topic)
    if normalized_topic:
        return normalized_topic
    try:
        rel = source.relative_to(workspace)
    except ValueError:
        return "unclassified"
    parts = rel.parts
    # New industry/ structure: industry/<slug>/companies/<ticker>/...
    if len(parts) >= 4 and parts[0] == "industry" and parts[2] == "companies":
        return f"industry/{parts[1]}/companies/{parts[3]}"
    if len(parts) >= 2 and parts[0] == "industry":
        return f"industry/{parts[1]}"
    # Legacy topics/ structure (still supported for backward compat)
    if len(parts) >= 5 and parts[0] == "topics" and parts[1] in TOPIC_NAMESPACES and parts[3] in TOPIC_OPERATION_DIRS:
        return f"{parts[1]}/{parts[2]}"
    if len(parts) >= 4 and parts[0] == "topics" and parts[2] in ("_raw", "_inbox", "_cache"):
        return parts[1]
    if len(parts) >= 3 and parts[0] == "_inbox" and parts[1] not in ("", "."):
        root_topic = "/".join(part for part in parts[1:-1] if part not in ("", "."))
        return root_topic or "unclassified"
    if len(parts) >= 2 and parts[0] == "_inbox":
        return "unclassified"
    return "unclassified"


def output_path_for(source: Path, workspace: Path, cache_root: Path | None, topic: str | None) -> Path:
    resolved = resolve_topic(source, workspace, topic)
    if cache_root:
        return cache_root / resolved / f"{source.stem}.md"
    return workspace / resolved / "_cache" / f"{source.stem}.md"


def read_cache_metadata(path: Path) -> dict[str, str]:
    if not path.exists():
        return {}
    text = path.read_text(encoding="utf-8", errors="replace")
    if not text.startswith("<!--"):
        return {}
    end = text.find("-->")
    if end == -1:
        return {}
    metadata: dict[str, str] = {}
    for line in text[4:end].splitlines():
        if ":" in line:
            key, value = line.split(":", 1)
            metadata[key.strip()] = value.strip()
    return metadata


def collision_safe_output_path(target: Path, source: Path) -> Path:
    suffix_token = source.suffix.lower().lstrip(".") or "file"
    return target.with_name(f"{source.stem}-{suffix_token}.md")


def candidate_files(source: Path, recursive: bool) -> list[Path]:
    if source.is_file():
        return [source]
    pattern = "**/*" if recursive else "*"
    return sorted(path for path in source.glob(pattern) if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS)


def _source_is_in_inbox(source: Path, workspace: Path) -> bool:
    try:
        source.resolve().relative_to((workspace / "_inbox").resolve())
        return True
    except ValueError:
        pass
    try:
        rel = source.resolve().relative_to(workspace.resolve())
        if len(rel.parts) >= 5 and rel.parts[0] == "topics" and rel.parts[1] in TOPIC_NAMESPACES and rel.parts[3] == "_inbox":
            return True
        if len(rel.parts) >= 3 and rel.parts[0] == "topics" and rel.parts[2] == "_inbox":
            return True
    except ValueError:
        pass
    return False


def _check_topic_exists(workspace: Path, topic: str) -> None:
    if topic == "unclassified":
        return
    topic_dir = workspace / topic
    if not topic_dir.is_dir():
        raise IngestError(
            f"Topic '{topic}' does not exist. Run new-session first to create the topic directory first.\n"
            f"Missing: {topic_dir}"
        )


def _move_to_raw(source: Path, workspace: Path, topic: str, category: str) -> Path:
    raw_dir = workspace / topic / "_raw" / category
    raw_dir.mkdir(parents=True, exist_ok=True)
    dest = raw_dir / source.name
    shutil.move(str(source), str(dest))
    return dest


def cache(source: Path, workspace: Path, cache_root: Path | None, topic: str | None, category: str | None, force: bool, extract_figures: bool = True) -> dict[str, Any]:
    profile = detect_format(source)
    target = output_path_for(source, workspace, cache_root, topic)
    resolved_topic = resolve_topic(source, workspace, topic)
    _check_topic_exists(workspace, resolved_topic)
    resolved_category = category or infer_category(source, profile)
    cache_parent = target.parent if not cache_root else None
    if target.exists() and not force:
        current_hash = sha256_file(source)
        metadata = read_cache_metadata(target)
        cached_source_path = metadata.get("source_path")
        cached_hash = metadata.get("source_sha256")
        if cached_source_path != str(source) or cached_hash != current_hash:
            target = collision_safe_output_path(target, source)
            if not target.exists():
                result = route_converter(source, profile, cache_parent=cache_parent, extract_figures=extract_figures)
                target.parent.mkdir(parents=True, exist_ok=True)
                target.write_text(markdown_header(source, result) + result.markdown.rstrip() + "\n", encoding="utf-8")
                raw_dest = None
                moved = False
                if _source_is_in_inbox(source, workspace):
                    raw_dest = str(_move_to_raw(source, workspace, resolved_topic, resolved_category))
                    moved = True
                entry: dict[str, Any] = {
                    "source": str(source),
                    "cache": str(target),
                    "status": "converted",
                    "converter": result.converter,
                    "precision": result.precision,
                    "precision_level": result.precision_level,
                    "document_type": result.document_type,
                    "route": f"{result.route}-collision-safe-cache",
                    "page_count": result.page_count,
                    "table_count": result.table_count,
                    "ocr_required": result.ocr_required,
                    "dependency_status": result.dependency_status,
                    "category": resolved_category,
                }
                if moved:
                    entry["moved_to_raw"] = raw_dest
                return entry

        entry = {
            "source": str(source),
            "cache": str(target),
            "status": "skipped",
            "converter": "cache-reuse",
            "precision": "existing cache; use --force to overwrite",
            "precision_level": "existing",
            "document_type": profile.document_type,
            "route": "cache-reuse",
            "page_count": profile.page_count,
            "table_count": profile.table_count,
            "ocr_required": profile.ocr_required,
        }
        if _source_is_in_inbox(source, workspace):
            entry["moved_to_raw"] = str(_move_to_raw(source, workspace, resolved_topic, resolved_category))
        return entry

    result = route_converter(source, profile, cache_parent=cache_parent, extract_figures=extract_figures)
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(markdown_header(source, result) + result.markdown.rstrip() + "\n", encoding="utf-8")
    raw_dest = None
    moved = False
    if _source_is_in_inbox(source, workspace):
        raw_dest = str(_move_to_raw(source, workspace, resolved_topic, resolved_category))
        moved = True
    entry = {
        "source": str(source),
        "cache": str(target),
        "status": "converted",
        "converter": result.converter,
        "precision": result.precision,
        "precision_level": result.precision_level,
        "document_type": result.document_type,
        "route": result.route,
        "page_count": result.page_count,
        "table_count": result.table_count,
        "ocr_required": result.ocr_required,
        "dependency_status": result.dependency_status,
        "category": resolved_category,
    }
    if moved:
        entry["moved_to_raw"] = raw_dest
    return entry


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Ingest raw research materials into the topic _cache/ directory markdown.")
    parser.add_argument("source_path", nargs="?", help="Source file or directory to ingest.")
    parser.add_argument("--workspace", help="Research workspace root. If omitted, discover from source path.")
    parser.add_argument("--cache-root", help="Override cache root. Defaults to <workspace>/<topic>/_cache.")
    parser.add_argument("--topic", help="Namespaced topic for organizing raw/cache (e.g. 'industry/space-launch' or 'company/rklb').")
    parser.add_argument("--category", help="Document category: filings, transcripts, sellside, industry, irdecks, datasets. Auto-inferred if omitted.")
    parser.add_argument("--bucket", help="Deprecated. Use --topic instead.")
    parser.add_argument("--force", action="store_true", help="Overwrite existing cache files.")
    parser.add_argument("--recursive", action="store_true", help="Recursively ingest supported files in a directory.")
    parser.add_argument("--check-deps", action="store_true", help="Print ingest dependency matrix as JSON and exit.")
    parser.add_argument("--extract-figures", action="store_true", default=True, dest="extract_figures",
                        help="Extract figures/charts from documents during ingest (default: True)")
    parser.add_argument("--no-extract-figures", action="store_false", dest="extract_figures",
                        help="Skip figure extraction")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    if args.check_deps:
        print(json.dumps(dependency_matrix(), ensure_ascii=False, indent=2))
        return 0

    if not args.source_path:
        print(json.dumps({"status": "failed", "error": "source_path is required unless --check-deps is used"}, ensure_ascii=False, indent=2))
        return 1

    source = Path(args.source_path).expanduser().resolve()
    if not source.exists():
        print(json.dumps({"status": "failed", "error": f"Source path does not exist: {source}"}, ensure_ascii=False, indent=2))
        return 1

    try:
        workspace = Path(args.workspace).expanduser().resolve() if args.workspace else discover_workspace(source)
        cache_root = Path(args.cache_root).expanduser().resolve() if args.cache_root else None
        topic = args.topic or args.bucket
        files = candidate_files(source, args.recursive)
        if not files:
            raise IngestError("No supported files found to ingest.")

        results = []
        failed = []
        for file_path in files:
            try:
                results.append(cache(file_path.resolve(), workspace, cache_root, topic, args.category, args.force, extract_figures=args.extract_figures))
            except Exception as exc:
                profile: dict[str, Any] = {}
                try:
                    detected = detect_format(file_path.resolve())
                    profile = {
                        "document_type": detected.document_type,
                        "page_count": detected.page_count,
                        "table_count": detected.table_count,
                        "ocr_required": detected.ocr_required,
                    }
                except Exception:
                    pass
                failed.append({"source": str(file_path), "status": "failed", "error": str(exc), **profile})

        payload = {
            "workspace": str(workspace),
            "converted": sum(1 for item in results if item["status"] == "converted"),
            "skipped": sum(1 for item in results if item["status"] == "skipped"),
            "failed": len(failed),
            "results": results,
            "errors": failed,
        }
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 1 if failed else 0
    except Exception as exc:
        print(json.dumps({"status": "failed", "error": str(exc)}, ensure_ascii=False, indent=2))
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
