#!/usr/bin/env python3
"""extract-pptx — PowerPoint to text/markdown. Lightweight.

Usage:
  python extract-pptx.py <file>              # plain text per slide
  python extract-pptx.py <file> --markdown    # markdown with slide titles
  python extract-pptx.py <file> --json        # structured JSON
"""
from __future__ import annotations

import argparse
import json
import sys


def _clean(text: str) -> str:
    return text.strip()


def extract(filepath: str) -> list[dict]:
    """Extract all slides. Returns [{number, title, text, notes, tables}]."""
    try:
        import pptx
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install --user python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = pptx.Presentation(filepath)
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        title = ""
        texts = []
        tables = []

        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == pptx.enum.shapes.PP_PLACEHOLDER.TITLE:
                title = _clean(shape.text)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = _clean(para.text)
                    if t:
                        texts.append(t)
            elif shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [_clean(cell.text) for cell in row.cells]
                    rows.append(cells)
                if rows:
                    tables.append({"headers": rows[0], "rows": rows[1:] if len(rows) > 1 else []})

        # Speaker notes
        notes = ""
        if slide.has_notes_slide:
            notes = _clean(slide.notes_slide.notes_text_frame.text)

        slides.append({
            "number": i,
            "title": title,
            "text": texts,
            "notes": notes,
            "tables": tables,
        })

    return slides


def to_markdown(slides: list[dict]) -> str:
    """Convert slides to markdown."""
    lines = []
    for s in slides:
        lines.append(f"## Slide {s['number']}")
        if s["title"]:
            lines.append(f"**{s['title']}**")
        lines.append("")
        for t in s["text"]:
            lines.append(t)
            lines.append("")
        for tbl in s["tables"]:
            if tbl["headers"]:
                lines.append("| " + " | ".join(tbl["headers"]) + " |")
                lines.append("|---" * len(tbl["headers"]) + "|")
                for row in tbl["rows"]:
                    lines.append("| " + " | ".join(row) + " |")
            lines.append("")
        if s["notes"]:
            lines.append(f"> Notes: {s['notes']}")
            lines.append("")
        lines.append("---")
        lines.append("")
    return "\n".join(lines)


def main():
    p = argparse.ArgumentParser(description="PPTX to text/markdown/JSON")
    p.add_argument("file", help="PPTX file path")
    p.add_argument("--markdown", action="store_true", help="Output as markdown")
    p.add_argument("--json", action="store_true", help="Output as structured JSON")
    args = p.parse_args()

    slides = extract(args.file)

    if args.json:
        print(json.dumps(slides, ensure_ascii=False, indent=2))
    elif args.markdown:
        print(to_markdown(slides))
    else:
        for s in slides:
            print(f"Slide {s['number']}: {s['title']}")
            if s["notes"]:
                print(f"  Notes: {s['notes']}")


if __name__ == "__main__":
    main()
